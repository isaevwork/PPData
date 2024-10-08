import os
import warnings
from PIL import Image
from pptx.dml.color import RGBColor
from openpyxl import load_workbook
from pptx import Presentation
from datetime import datetime
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT


def draw_table(text):
    # Определение длины самой длинной строки в тексте
    max_length = max(len(line) for line in text.split('\n'))

    # Верхняя граница таблицы
    print("-" * (max_length + 4))

    # Выводим текст с отступами слева и справа
    for line in text.split('\n'):
        print(f"| {line.center(max_length)} |")

    # Нижняя граница таблицы
    print("-" * (max_length + 4))


# Текст для помещения внутрь таблицы
text = """Предупреждение: использование этого скрипта полностью лежит на ответственности конечного пользователя.
Автор скрипта не несет никакой ответственности!
Мы просим всех пользователей внимательно изучить результаты этого скрипта перед использованием."""


def main():
    work_folder = os.path.join(os.environ['USERPROFILE'], 'Downloads', 'WORK')
    os.chdir(work_folder)

    # Вызываем функцию и передаем текст внутрь таблицы
    draw_table(text)

    extensions = [".jpg", ".png", ".jpeg", ".gif"]  # Расширения изображений для переименования

    for folder_name in os.listdir():
        if os.path.isdir(folder_name):
            os.chdir(folder_name)
            renamed_files = set()  # Сет для отслеживания уже переименованных файлов
            for filename in os.listdir():
                if any(filename.lower().endswith(ext) for ext in extensions):
                    if folder_name not in filename and filename not in renamed_files:
                        new_filename = f"{folder_name}_{filename}"
                        os.rename(filename, new_filename)
                        renamed_files.add(new_filename)  # Добавляем новое имя файла в сет
            os.chdir('..')


main()

# Путь к папке work
user_profile = os.getenv('USERPROFILE')
work_folder = os.path.join(user_profile, 'Downloads', 'work')

# Список для хранения путей к файлам Excel
excel_files = []

# Глобальная переменная для хранения имени главной папки
folder_name = None

# Глобальная переменная для хранения Excel файла
wb = None

# Путь к файлу Excel с таблицей
# Рекурсивно обходим все подпапки внутри папки work
for root, dirs, files in os.walk(work_folder):
    # Просматриваем все файлы в текущей подпапке
    for file in files:
        # Составляем полный путь к текущему файлу
        file_path = os.path.join(root, file)
        # Если файл имеет расширение .xlsx и находится не в корневой папке work, добавляем его путь в список
        if file.endswith(".xlsx") and root != work_folder:
            excel_files.append(file_path)

# Если найден хотя бы один файл Excel, загружаем первый из них
if excel_files:
    excel_file_path = excel_files[0]
    with warnings.catch_warnings():
        warnings.simplefilter("ignore")
        wb = load_workbook(filename=excel_file_path, data_only=True)
        print(f"Excel файл {os.path.basename(excel_file_path)} найден в папке {os.path.dirname(excel_file_path)}.")

        # Извлекаем имя папки из пути
        folder_name = os.path.basename(os.path.dirname(excel_file_path))
else:
    print(f" Файл Excel не найден в подпапках папки work.")

# Загружаем презентацию
prs = Presentation(os.path.join(os.getenv('USERPROFILE'), 'Downloads', 'work', 'FDTemp.pptx'))
image_folder = os.path.join(work_folder, folder_name)
# -----------------------------------------------------------------------------------------
# Определение параметров текстового блока
left = Inches(6.8)
top = Inches(7.8)
width = Inches(3.3)
height = Inches(2)

# Список индексов слайдов, для которых нужно создать текстовые блоки с ФИО
slide_indexes = [1, 2, 3, 4, 5, 6, 7, 8, 9, 10, 11, 12, 13, 14, 15, 16, 17, 18, 19, 20, 21, 22]

# Цикл для создания текстовых блоков на каждом слайде
for index in slide_indexes:
    name_textbox = prs.slides[index].shapes.add_textbox(left, top, width, height)
    name_textbox.rotation = 270
    tf = name_textbox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = f"{folder_name}"
    p.font.bold = False
    p.font.size = Pt(12)
    p.font.name = "Montserrat"


def insert_images(names, positions, idx):
    """
    Функция для добавления изображений на конкретный слайд презентации.
    Args:
        names (list): Список имен изображений.
        positions (list): Словарь с позициями изображений для каждого слайда.
        idx (int): Индекс слайда, на который добавляются изображения.
    """
    extensions = [".jpg", ".png", ".jpeg", ".gif"]  # Расширения изображений для проверки
    slide = prs.slides[idx]  # Получаем слайд по индексу

    for name, position in zip(names, positions):
        for extension in extensions:
            image_path = os.path.join(image_folder, name + extension)
            if os.path.exists(image_path):
                img_left, img_top, img_width, img_height = position
                slide.shapes.add_picture(image_path, img_left, img_top, img_width, img_height)
                break
        else:
            print(f" Изображение {name} не найдено на слайде {idx}.")


def format_with_comma(number):
    """
    Форматирует число с одним десятичным знаком и заменяет точку на запятую.
    """
    formatted_number = "{:,.1f}".format(number)
    return formatted_number.replace(".", ",")


def process_string(input_string, round_num):
    # Преобразует строку с плавающей запятой в формате "12,8" в число
    # заменяет запятую на точку, при необходимости округляет до двух знаков после запятой
    # возвращает строку с плавающей запятой в формате "12,80"
    if isinstance(input_string, (int, float)):
        # Если входная строка уже является числом, преобразуем ее к типу float
        number = float(input_string)
    else:
        if ',' in input_string:
            # Заменяем запятую на точку
            input_string = input_string.replace(',', '.')

        try:
            number = float(input_string)
        except ValueError:
            return "Некорректный формат числа"  # Возвращаем ошибку, если не удается преобразовать в число

    # Убираем знак минус посредством abs(), если он есть
    number = abs(number)

    # Округляем число до двух знаков после запятой
    number = round(number, round_num)

    # Заменяем точку на запятую и добавляем 0 до двух знаков после запятой
    result_string = f"{number:.{round_num}f}".replace('.', ',')

    return result_string


def format_float_with_zeros(num):
    """
    Форматирует число с плавающей точкой с нулями в конце.
    """
    # Проверяем, является ли число целым
    if num == int(num):
        str_num = str(num)
        if '.' not in str_num:
            return f"{num}.0"  # Добавляем точку и ноль в конце, если нет десятичной части
        else:
            return str_num  # Возвращаем число без изменений, если десятичная часть уже присутствует
    else:
        return str(num)  # Возвращаем число без изменений


def get_text_color(last_value):
    """
    Возвращает цвет текста в зависимости от значения value.
    """
    if isinstance(last_value, (int, float)) and last_value is not None:
        if last_value <= -3.1:
            return RGBColor(255, 0, 0)  # Красный
        elif -3 <= last_value <= -2.1:
            return RGBColor(0, 0, 255)  # Синий
        elif -2 <= last_value <= -1.1:
            return RGBColor(6, 102, 6)  # Зеленый
        elif -1 <= last_value <= 1:
            return RGBColor(0, 0, 0)  # Черный
        elif 1.1 <= last_value <= 2:
            return RGBColor(6, 102, 6)  # Зеленый
        elif 2.1 <= last_value <= 3:
            return RGBColor(0, 0, 255)  # Синий
        elif last_value >= 3.1:
            return RGBColor(255, 0, 0)  # Красный
    return RGBColor(0, 0, 0)  # Черный цвет по умолчанию


def add_text_to_slide(presentation, slide_index, slide_data, current_left, current_top, cell_width, cell_height,
                      font_s):
    for i, row_data in enumerate(slide_data):
        for j, value in enumerate(row_data):
            # Рассчитываем координаты для текущей ячейки
            cell_left = current_left + j * cell_width
            cell_top = current_top + i * cell_height

            # Получаем значение last_value из последней ячейки текущей строки
            last_value = row_data[-1]

            # Получаем цвет текста на основе значения last_value
            color = get_text_color(last_value)

            # Преобразуем значение в строку, если оно не None
            if value is not None:
                if isinstance(value, (int, float)):
                    text_value = format_float_with_zeros(round(value, 1))
                else:
                    text_value = str(value)
            else:
                text_value = ""

            if j == 1:
                cell_left += Inches(0.8)
            if j == 2:
                cell_left += Inches(0.8)
            if j == 3:
                cell_left += Inches(0.8)
            if j == 4:
                cell_left += Inches(0.8)

            if i == 6:
                cell_top += Inches(0.04)
            if i == 7:
                cell_top += Inches(0.04)

            # Добавление текстового блока на слайд с текущими координатами и цветом
            table_frame = presentation.slides[slide_index].shapes.add_textbox(cell_left, cell_top, cell_width,
                                                                              cell_height).text_frame
            q = table_frame.add_paragraph()
            q.text = text_value
            q.font.size = font_s
            q.font.name = "Montserrat"
            q.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            q.font.color.rgb = color  # Устанавливаем цвет текста


def crop_and_resize_image(input_image_path, output_image_path):
    with Image.open(input_image_path) as img:
        # Получаем оригинальные размеры
        original_width, original_height = img.size

        # Проверяем ширину
        if original_width < 1000:
            print(f"Изображение {input_image_path} имеет ширину меньше 1000px, сохранение без изменений.")
            img.save(output_image_path)  # Сохраняем изображение без изменений
            return

        # Вычисляем новые размеры
        new_height = original_height  # Высота остаётся прежней
        new_width = new_height - 150  # Ширина на 150 пикселей меньше высоты

        # Обрезка по центру
        if new_width > 0:
            left = (original_width - new_width) // 2
            top = 0
            right = left + new_width
            bottom = original_height

            # Обрезаем изображение
            img_cropped = img.crop((left, top, right, bottom))

            # Сохраняем итоговое изображение
            img_cropped.save(output_image_path)
        else:
            print("Ошибка: новое значение ширины меньше или равно нулю.")


def rename_image(old_name, new_name):
    temp_folder = os.path.join(image_folder, "temp")
    name_folder_element = os.path.basename(image_folder)

    # Проверяем существование папки temp, если ее нет, создаем
    if not os.path.exists(temp_folder):
        try:
            os.makedirs(temp_folder)
        except Exception as e:
            print(f"Ошибка при создании папки temp: {str(e)}")
            return

    for extension in ['.jpg', '.png']:
        img_path = os.path.join(image_folder, f"{old_name}{extension}")
        if os.path.exists(img_path):
            break
    else:
        print(f" Изображение {old_name} не найдено.")
        return

    try:
        # Открываем изображение
        image = Image.open(img_path)
        # Формируем новый путь для сохранения в папку temp
        new_img_path = os.path.join(temp_folder, f"{name_folder_element}_{new_name}{extension}")
        # Создаем копию изображения с новым именем
        image_copy = image.copy()
        # Сохраняем копию в новом пути
        image_copy.save(new_img_path)
        print(
            f""" Изображение "{old_name}" успешно скопировано с новым именем "{name_folder_element}_{new_name}" и сохранено в папку temp.""")
    except Exception as e:
        print(f" Ошибка при копировании и сохранении изображения: {str(e)}")


print("<-------------------------------------------------------------------------------------------------------->")

# Задаем имя пациента, врача и дату
left = Inches(2.9)  # Расстояние от правого края слайда
top = Inches(7.75)  # Расстояние от верхнего края слайда
width = Inches(4)  # Ширина, чтобы занять всю ширину слайда
height = Inches(0.5)  # Высота, чтобы занять всю высоту слайда
name_textbox = prs.slides[0].shapes.add_textbox(left, top, width, height)
tf = name_textbox.text_frame
tf.word_wrap = True
p = tf.add_paragraph()
p.text = f"{folder_name}"
p.font.size = Pt(14)
p.font.name = "Montserrat Medium"

date_left = Inches(5)
date_top = Inches(8.9)
date_width = Inches(4)
name_textbox = prs.slides[0].shapes.add_textbox(date_left, date_top, date_width, height)
tf_date = name_textbox.text_frame
tf_date.word_wrap = True
p_date = tf_date.add_paragraph()
p_date.text = f"{datetime.today().strftime('%d.%m.%Y')}"
print(f"Слайд 1 сформирован")

print("<-------------------------------------------------------------------------------------------------------->")
# Слайд 2
ws2 = wb["Лист2"]
face_width = ws2['B10'].value  # Ширина
face_height = ws2['B11'].value  # Высота
facial_index = ws2['B12'].value  # Лицевой индекс
nosocomial_angle = ws2['B13'].value  # Носолицевой угол
nasal_angle = ws2['B14'].value  # Носоподбородочный угол
labial_angle = ws2['B15'].value  # Носогубный угол
chin_facial_angle = ws2['B16'].value  # Подбородочно-лицевой угол
soft_tissues_angle = ws2['B17'].value  # Угол выпуклости мягких тканей лица
upper_lip_position = ws2['B18'].value  # Положение верхней губы
lower_lip_position = ws2['B19'].value  # Положение нижней губы
profile_type = ''

if upper_lip_position < -6:
    upper_lip_position_status = 'ретроположение'
elif -6 <= upper_lip_position <= -2:
    upper_lip_position_status = 'в норме'
else:  # upper_lip_position > -2
    upper_lip_position_status = 'антеположение'

if lower_lip_position < -4:
    lower_lip_position_status = 'ретроположение'
elif -4 <= lower_lip_position <= 0:
    lower_lip_position_status = 'в норме'
else:  # lower_lip_position > 0
    lower_lip_position_status = 'антеположение'

if soft_tissues_angle < 163:
    profile_type = 'выпуклый'
elif 163 <= soft_tissues_angle <= 175:
    profile_type = 'в пределах нормы'
elif 176 <= soft_tissues_angle <= 180:
    profile_type = 'прямой'
elif soft_tissues_angle > 180:
    profile_type = 'вогнутый'


def add_num_to_slide(prs_n, slide_index, left_n, top_n, text_n, width_n=Inches(2.5), f_color=RGBColor(0, 0, 0),
                     font_size_n=10, font_name="Montserrat", bold=False):
    slide = prs_n.slides[slide_index]
    textbox = slide.shapes.add_textbox(left_n, top_n, width_n, Inches(0))
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = text_n
    p.font.bold = bold
    p.font.size = Pt(font_size_n)
    p.font.color.rgb = f_color
    p.font.name = font_name


add_num_to_slide(prs, 2, Inches(4.5), Inches(1.81), f"{format_with_comma(face_width)}")
add_num_to_slide(prs, 2, Inches(4.65), Inches(2.08), f"{format_with_comma(face_height)}")
add_num_to_slide(prs, 2, Inches(4.36), Inches(2.315), f"{format_with_comma(facial_index)}")
add_num_to_slide(prs, 2, Inches(1.9), Inches(5.39), f"{format_with_comma(nosocomial_angle)}°")
add_num_to_slide(prs, 2, Inches(2.55), Inches(5.65), f"{format_with_comma(nasal_angle)}°")
add_num_to_slide(prs, 2, Inches(1.87), Inches(6.15), f"{format_with_comma(labial_angle)}°")
add_num_to_slide(prs, 2, Inches(2.8), Inches(6.4), f"{format_with_comma(chin_facial_angle)}°")
add_num_to_slide(prs, 2, Inches(1.2), Inches(7.15), f"{format_with_comma(soft_tissues_angle)}°")

add_num_to_slide(prs, 2, Inches(5.25), Inches(5.4), f"{profile_type}")  # Расположение Типа профиля

add_num_to_slide(prs, 2, Inches(6.1), Inches(5.65), f"{upper_lip_position_status} —")  # Положение верхней губы
add_num_to_slide(prs, 2, Inches(4.15), Inches(5.85),
                 f"{format_with_comma(upper_lip_position)} мм (N = -4,0 мм ±2,0 мм)")

add_num_to_slide(prs, 2, Inches(6.1), Inches(6.15), f"{lower_lip_position_status} —")  # Положение нижней губы
add_num_to_slide(prs, 2, Inches(4.15), Inches(6.35),
                 f"{format_with_comma(lower_lip_position)} мм (N = -2,0 мм ±2,0 мм)")

# Массив имен изображений с префиксом папки

images_name_2 = [f"{folder_name}_{image}" for image in ["a", "s", "d", "f"]]

images_position_2 = [
    (Inches(0.4), Inches(1.5), Inches(2.6), Inches(3.6)),
    (Inches(5.4), Inches(1.5), Inches(2.6), Inches(3.6)),
    (Inches(0.7), Inches(7.8), Inches(3), Inches(3.6)),
    (Inches(4.6), Inches(7.7), Inches(3), Inches(3.7))
]

# Проходим по каждому имени изображения в images_name_2
for image_name in images_name_2:
    # Проверяем наличие .jpg и .png файлов, создаем полный путь
    for ext in ['.jpg', '.png']:
        input_image_path = os.path.join(image_folder, image_name + ext)

        # Проверяем, существует ли файл
        if os.path.exists(input_image_path):
            output_image_path = os.path.join(image_folder, image_name + ext)
            # Вызываем функцию для обрезки и изменения размера
            crop_and_resize_image(input_image_path, output_image_path)
            break  # Выход из цикла, как только файл найден

insert_images(images_name_2, images_position_2, 2)
print(f"Слайд 2 сформирован")

print("<-------------------------------------------------------------------------------------------------------->")
# Слайд 3 новый
images_name_3 = [f"{folder_name}_{image}" for image in ["aa"]]
images_position_3 = [
    (Inches(0.6), Inches(1.4), Inches(7.1), Inches(8.2)),
]

insert_images(images_name_3, images_position_3, 3)
print(f"Слайд 3 сформирован")

print("<-------------------------------------------------------------------------------------------------------->")
# Слайд 4 новый
images_name_4 = [f"{folder_name}_{image}" for image in ["ss"]]
images_position_4 = [
    (Inches(0.6), Inches(1.4), Inches(7.1), Inches(8.2)),
]

insert_images(images_name_4, images_position_4, 4)
print(f"Слайд 4 сформирован")

print("<-------------------------------------------------------------------------------------------------------->")
# Слайд 5 новый
images_name_5 = [f"{folder_name}_{image}" for image in ["dd"]]
images_position_5 = [
    (Inches(0.6), Inches(1.4), Inches(7.1), Inches(8.2)),
]

insert_images(images_name_5, images_position_5, 5)
print(f"Слайд 5 сформирован")

print("<-------------------------------------------------------------------------------------------------------->")
# Слайд 6 новый
slideThree_MT = list(ws2.iter_rows(min_row=2, max_row=9, min_col=16, max_col=17, values_only=True))

# Размер и положение данных на слайде
c_left = Inches(2.8)  # Левая граница
c_top = Inches(5.53)  # Верхняя граница
c_width = Inches(2.5)  # Ширина ячейки
c_height = Inches(0.275)  # Высота ячейки
f_size = Pt(12)  # Размер шрифта


def transform_data(data):
    transformed_data = []
    for idx, sublist in enumerate(data):
        transformed_sublist = []
        for item in sublist:
            if isinstance(item, (int, float)):
                # Проверяем, является ли текущий подсписок последним или предпоследним в массиве данных
                if idx == len(data) - 2 or idx == len(data) - 1:
                    transformed_sublist.append('{:.1f}%'.format(item * 100).replace('.', ','))
                else:
                    transformed_sublist.append(round(item, 2))
            else:
                transformed_sublist.append(item)
        transformed_data.append(transformed_sublist)
    return transformed_data


transformed_dataframe = transform_data(slideThree_MT)

# Размещение данных на слайде
for i, row_data in enumerate(transformed_dataframe):
    for j, value in enumerate(row_data):
        # Рассчитываем координаты для текущей ячейки
        cell_left = c_left + j * c_width
        cell_top = c_top + i * c_height

        # Добавление текстового блока на слайд с текущими координатами
        text_frame = prs.slides[6].shapes.add_textbox(cell_left, cell_top, c_width, c_height).text_frame
        p = text_frame.add_paragraph()
        p.text = str(value)
        p.font.size = f_size
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        p.font.name = "Montserrat"


def fill_table(present, slide_index, slide_data, cl, ct, cw, ch, fs, column_offsets):
    for i, row_datas in enumerate(slide_data):
        for j, ce_value in enumerate(row_datas):
            # Рассчитываем координаты для текущей ячейки
            ce_left = cl + j * cw

            # Преобразуем значение в строку, если оно не None
            text_value = str(round(ce_value, 2)) if isinstance(ce_value, (int, float)) else str(ce_value)

            # Получаем информацию о смещении для текущего столбца
            offset_info = column_offsets.get(j)
            if offset_info:
                add_offset, offset_value = offset_info
                if add_offset:
                    ce_left += offset_value
                else:
                    ce_left = offset_value

            # Добавление текстового блока на слайд с текущими координатами и цветом
            table_frame = present.slides[slide_index].shapes.add_textbox(ce_left, ct, cw, ch).text_frame
            p = table_frame.add_paragraph()
            p.text = text_value
            p.font.size = fs
            p.font.name = "Montserrat"
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER


up_data3 = list(ws2.iter_rows(min_row=2, max_row=2, min_col=2, max_col=13, values_only=True))
lower_data3 = list(ws2.iter_rows(min_row=3, max_row=3, min_col=2, max_col=13, values_only=True))

column_offsets_up = {
    1: (True, Inches(0.1)),
    2: (True, Inches(0.05)),
    7: (False, Inches(4.65)),
    8: (False, Inches(5.15)),
    9: (False, Inches(5.65)),
    10: (False, Inches(6.2)),
    11: (False, Inches(6.76))
}
column_offsets_lower = {
    1: (True, Inches(0.15)),
    2: (True, Inches(0.2)),
    3: (True, Inches(0.2)),
    4: (False, Inches(3.3)),
    5: (False, Inches(3.74)),
    6: (False, Inches(4.1)),
    10: (False, Inches(5.9)),
    11: (False, Inches(6.5))
}

fill_table(prs, 6, up_data3, Inches(0.9), Inches(2.45), Inches(0.55), Inches(0.27), Pt(14), column_offsets_up)
fill_table(prs, 6, lower_data3, Inches(1.3), Inches(3.65), Inches(0.45), Inches(0.27), Pt(14), column_offsets_lower)
print(f"Слайд 6 сформирован")

print("<-------------------------------------------------------------------------------------------------------->")
# Слайд 7
images_name_7 = [f"{folder_name}_{image}" for image in ["2q", "3q", "2w", "3w", "00", "1r"]]
images_position_7 = [
    (Inches(0.9), Inches(1.2), Inches(2.9), Inches(2.75)),
    (Inches(4.4), Inches(1.2), Inches(2.9), Inches(2.75)),
    (Inches(0.9), Inches(4), Inches(2.9), Inches(2.75)),
    (Inches(4.4), Inches(4), Inches(2.9), Inches(2.75)),
    (Inches(0.8), Inches(9.1), Inches(2.4), Inches(2.2)),
    (Inches(4.3), Inches(9.3), Inches(3.2), Inches(2))
]

insert_images(images_name_7, images_position_7, 7)
rename_image(images_name_7[4], "размеры апикальных базисов вч и нч")

upper_dentition_rows = list(ws2.iter_rows(min_row=23, max_row=26, min_col=2, max_col=4, values_only=True))
lower_dentition_rows = list(ws2.iter_rows(min_row=23, max_row=26, min_col=5, max_col=7, values_only=True))


def convert_list(data):
    transformed_data = []
    for sublist in data:
        transformed_sublist = []
        for item in sublist:
            transformed_sublist.append(item)
        transformed_data.append(transformed_sublist)
    return transformed_data


# Преобразование данных и размещение на слайде
transformed_upper_rows = convert_list(upper_dentition_rows)
transformed_lower_rows = convert_list(lower_dentition_rows)


def place_data_on_slide(data_rows, left_mar, top_mar, cell_width, cell_height, size_f, slide_index):
    for i, row_data in enumerate(data_rows):
        for j, value in enumerate(row_data):
            cell_left = left_mar + j * cell_width
            cell_top = top_mar + i * cell_height

            new_text_frame = prs.slides[slide_index].shapes.add_textbox(cell_left, cell_top, cell_width,
                                                                        cell_height).text_frame
            p = new_text_frame.add_paragraph()
            p.text = format_with_comma(value)
            p.font.size = Pt(size_f)
            p.font.name = "Montserrat"
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER


# Задание параметров для верхних и нижних данных
upper_left_mar = Inches(2.85)
lower_left_mar = Inches(5.3)
top_mar = Inches(7.15)

width = Inches(0.74)
height = Inches(0.34)
font_size = 13

# Размещение данных на слайде для верхних и нижних строк
place_data_on_slide(transformed_upper_rows, upper_left_mar, top_mar, width, height, font_size, 7)
place_data_on_slide(transformed_lower_rows, lower_left_mar, top_mar, width, height, font_size, 7)
print(f" Слайд №7 сформирован")

print("<-------------------------------------------------------------------------------------------------------->")
# Слайд 8
# Массив имен изображений с префиксом папки
images_name_8 = [f"{folder_name}_{image}" for image in ["2e", "3e", "2r", "3y", "3r", "3t"]]
depth_shpee_right_curve = ws2['B29'].value
depth_shpee_left_curve = ws2['C29'].value

images_position_8 = [
    (Inches(0.6), Inches(1.2), Inches(3.3), Inches(3)),
    (Inches(4.4), Inches(1.2), Inches(3.3), Inches(3)),

    (Inches(0.6), Inches(5), Inches(3.1), Inches(3)),
    (Inches(4.5), Inches(5), Inches(3.1), Inches(3)),

    (Inches(0.6), Inches(9.2), Inches(3.4), Inches(1.4)),
    (Inches(4.2), Inches(9.2), Inches(3.4), Inches(1.4))
]
slide_index_8 = 8
insert_images(images_name_8, images_position_8, slide_index_8)

add_num_to_slide(prs, slide_index_8, Inches(1.4), Inches(11),
                 f"R = {process_string(depth_shpee_right_curve, 1)} мм (N = 1,5 мм)", Inches(3), RGBColor(0, 0, 0), 14)
add_num_to_slide(prs, slide_index_8, Inches(5.1), Inches(11),
                 f"L = {process_string(depth_shpee_left_curve, 1)} мм (N = 1,5 мм)", Inches(3), RGBColor(0, 0, 0), 14)
print(f" Слайд 8 сформирован")

print("<-------------------------------------------------------------------------------------------------------->")
# Слайд 9
images_name_9 = [f"{folder_name}_{image}" for image in ["1w", "1q", "1e", "1ww", "1qq", "1ee"]]
images_position_9 = [
    (Inches(0.8), Inches(1.3), Inches(3.2), Inches(1.8)),

    (Inches(0.8), Inches(3.3), Inches(3.2), Inches(1.8)),

    (Inches(0.8), Inches(5.2), Inches(3.2), Inches(1.8)),

    (Inches(4.4), Inches(1.3), Inches(3.1), Inches(1.8)),
    (Inches(4.4), Inches(3.3), Inches(3.1), Inches(1.8)),
    (Inches(4.4), Inches(5.2), Inches(3.1), Inches(1.8))
]
slide_index_9 = 9
insert_images(images_name_9, images_position_9, slide_index_9)
print(f" Слайд 9 сформирован")

print("<-------------------------------------------------------------------------------------------------------->")
# Слайд 10
# Массив имен изображений с префиксом папки
images_name_10 = [f"{folder_name}_{image}" for image in ["9", "6"]]
images_position_10 = [
    (Inches(0.6), Inches(1.5), Inches(7), Inches(4)),
    (Inches(0.6), Inches(6.3), Inches(7), Inches(4.6)),
]

insert_images(images_name_10, images_position_10, 10)
rename_image(images_name_10[0], "мягкие ткани")
rename_image(images_name_10[1], "костная ткань")
print(f" Слайд №10 сформирован")

print("<-------------------------------------------------------------------------------------------------------->")
# Слайд 11
images_name_11 = [f"{folder_name}_{image}" for image in ["11", "вч", "нч"]]
images_position_11 = [
    (Inches(0.5), Inches(1.1), Inches(7.4), Inches(3.4)),
    (Inches(0.5), Inches(5.1), Inches(7.2), Inches(2.9)),
    (Inches(0.5), Inches(8.5), Inches(7.2), Inches(2.9)),
]

insert_images(images_name_11, images_position_11, 11)
rename_image(images_name_11[0], "ОПТГ")
print(f" Слайд №11 сформирован")

print("<-------------------------------------------------------------------------------------------------------->")
# Слайд 12
# Массив имен изображений с префиксом папки
images_name_12 = [f"{folder_name}_{image}" for image in ["22", "1", "2"]]
images_position_12 = [
    (Inches(0.4), Inches(1.8), Inches(7.6), Inches(3.2)),
    (Inches(1.5), Inches(5.8), Inches(5.9), Inches(2.6)),
    (Inches(1.5), Inches(8.7), Inches(5.9), Inches(2.6)),
]
insert_images(images_name_12, images_position_12, 12)
rename_image(images_name_12[0], "ВНЧС")
rename_image(images_name_12[1], "ВНЧС прав")
rename_image(images_name_12[2], "ВНЧС лев")
print(f" Слайд 12 сформирован")

print("<-------------------------------------------------------------------------------------------------------->")
# Слайд 13
# Размеры и положения областей для изображений
images_name_13_444 = [f"{folder_name}_{image}" for image in ["444"]]
images_name_13 = [f"{folder_name}_{image}" for image in ["33", "44"]]
img_name13_1 = os.path.join(images_name_13_444[0] + ".jpg")

images_position_13 = [
    (Inches(0.5), Inches(7.5), Inches(3.5), Inches(3.5)),
    (Inches(4.1), Inches(7.5), Inches(3.5), Inches(3.5)),
]

prs.slides[13].shapes.add_picture(os.path.join(image_folder, images_name_13_444[0] + ".jpg"), Inches(1.2), Inches(1.4),
                                  Inches(6), Inches(5.5))
insert_images(images_name_13, images_position_13, 13)
rename_image(images_name_13[0], "ТРГ фронт")
rename_image(images_name_13[1], "SMV")
rename_image(images_name_13_444[0], "симметрия")
print(f" Слайд №13 сформирован")

print("<-------------------------------------------------------------------------------------------------------->")
# Слайд 14
images_name_14 = [f"{folder_name}_{image}" for image in ["4", "5", "55"]]
images_position_14 = [
    (Inches(0.4), Inches(1.4), Inches(3.6), Inches(3.6)),
    (Inches(4.2), Inches(1.4), Inches(3.6), Inches(3.6)),

    (Inches(1), Inches(5.8), Inches(6.2), Inches(5.6)),
]

insert_images(images_name_14, images_position_14, 14)
rename_image(images_name_14[0], "ТРГ прав")
rename_image(images_name_14[1], "ТРГ лев")
rename_image(images_name_14[2], "трассированная трг")
print(f" Слайд 14 сформирован")

print("<-------------------------------------------------------------------------------------------------------->")
# Слайд 15
images_name_15 = [f"{folder_name}_{image}" for image in ["56"]]
images_position_15 = [
    (Inches(0.65), Inches(1.6), Inches(7.1), Inches(8.5)),
]
slide_index_15 = 15
insert_images(images_name_15, images_position_15, slide_index_15)
print(f" Слайд 15 сформирован")

print("<-------------------------------------------------------------------------------------------------------->")
# Слайд 16
images_name_16 = [f"{folder_name}_{image}" for image in ["57"]]
images_position_16 = [
    (Inches(0.6), Inches(1.8), Inches(6.8), Inches(6.6)),
]
slide_index_16 = 16
insert_images(images_name_16, images_position_16, slide_index_16)

ws1 = wb["Лист1"]
params16_data = []

# Заполняем DataFrame данными из листа Excel
for row in ws1.iter_rows(min_row=29, max_row=36, min_col=2, max_col=6, values_only=True):
    params16_data.append(list(row))

# Размеры и положение данных на слайде

params16_left = Inches(1.3)  # Левая граница
params16_top = Inches(9.08)  # Верхняя граница
params16_width = Inches(1.12)  # Ширина ячейки
params16_height = Inches(0.27)  # Высота ячейки
font_size = Pt(11)  # Размер шрифта

add_text_to_slide(prs, 16, params16_data, params16_left, params16_top, params16_width, params16_height, font_size)
print(f" Слайд 16 сформирован")

print("<-------------------------------------------------------------------------------------------------------->")
# Слайд 17
images_name_17 = [f"{folder_name}_{image}" for image in ["77", "88"]]
images_position_17 = [
    (Inches(0.6), Inches(8.1), Inches(3.4), Inches(3.3)),
    (Inches(4.2), Inches(8.1), Inches(3.4), Inches(3.3)),
]

insert_images(images_name_17, images_position_17, 17)
rename_image([f"{folder_name}_{image}" for image in ["3"]][0], "фронтальный расчет")
rename_image(images_name_17[1], "Yaw ротация нижней челюсти")
up_params17_data = []
lower_params17_data = []

# Заполняем DataFrame данными из листа Excel
for row in ws1.iter_rows(min_row=9, max_row=14, min_col=2, max_col=6, values_only=True):
    up_params17_data.append(list(row))
for row in ws1.iter_rows(min_row=17, max_row=26, min_col=2, max_col=6, values_only=True):
    lower_params17_data.append(list(row))

# Размеры и положение данных на слайде
up_params17_left = Inches(1.2)  # Левая граница
up_params17_top = Inches(1.96)  # Верхняя граница
up_params17_width = Inches(1.12)  # Ширина ячейки
up_params17_height = Inches(0.27)  # Высота ячейки

lower_params17_left = Inches(1.2)  # Левая граница
lower_params17_top = Inches(4.3)  # Верхняя граница
lower_params17_width = Inches(1.12)  # Ширина ячейки
lower_params17_height = Inches(0.27)  # Высота ячейки

font_size = Pt(9)

add_text_to_slide(prs, 17, up_params17_data, up_params17_left, up_params17_top, up_params17_width, up_params17_height,
                  font_size)
add_text_to_slide(prs, 17, lower_params17_data, lower_params17_left, lower_params17_top, lower_params17_width,
                  lower_params17_height, font_size)
print(f" Слайд №17 сформирован")

print("<-------------------------------------------------------------------------------------------------------->")
# Слайд 18
images_name_18 = [f"{folder_name}_{image}" for image in ["333", "99"]]
images_position_18 = [
    (Inches(0.6), Inches(1.5), Inches(7), Inches(4.7)),
    (Inches(0.6), Inches(7), Inches(7), Inches(4.4)),
]

insert_images(images_name_18, images_position_18, 18)
rename_image(images_name_18[0], "аксиальные срезы")
print(f" Слайд №18 сформирован")

print("<-------------------------------------------------------------------------------------------------------->")
# Слайд 19
# Массив имен изображений с префиксом папки
images_name_19 = [f"{folder_name}_{image}" for image in ["222", "0"]]
images_position_19 = [
    (Inches(0.6), Inches(1.2), Inches(7), Inches(4)),
    (Inches(2), Inches(7.65), Inches(4.3), Inches(3.9)),
]

insert_images(images_name_19, images_position_19, 19)
rename_image(images_name_19[0], "корональные срезы")
rename_image(images_name_19[1], "воздухоносные пути")

# !!!!

airway_volume = list(ws1.iter_rows(min_row=26, max_row=27, min_col=13, max_col=15, values_only=True))

# Преобразование данных и размещение на слайде
transformed_airway_volume = convert_list(airway_volume)


def land_on_slide(data_rows, left_mar, top_mar, cell_width, cell_height, size_f, slide_index):
    for i, row_data in enumerate(data_rows):
        for j, value in enumerate(row_data):
            cell_left = left_mar + j * cell_width
            cell_top = top_mar + i * cell_height

            new_text_frame = prs.slides[slide_index].shapes.add_textbox(cell_left, cell_top, cell_width,
                                                                        cell_height).text_frame
            p = new_text_frame.add_paragraph()
            p.text = str(value)
            p.font.size = Pt(size_f)
            p.font.name = "Montserrat"
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER


left_margin = Inches(4.4)
top_mar = Inches(6.3)

width = Inches(1.1)
height = Inches(0.5)
font_size = 13

land_on_slide(transformed_airway_volume, left_margin, top_mar, width, height, font_size, 19)

print(f" Слайд №19 сформирован")

print("<-------------------------------------------------------------------------------------------------------->")
# Слайд 20
# Определяем тенденцию к классу
anb_value = ws1['L42'].value
beta_angle = ws1['L43'].value
wits_appraisal = ws1['L44'].value
sassouni = ws1['L125'].value
apdi_value = ws1['L46'].value
pnsa_value = ws1['C9'].value
jj_value = ws1['C10'].value
sna_value = ws1['C13'].value
ppsn_value = ws1['C14'].value
ton_index = ws2['P6'].value
general_Bolton_Index = ws2['P8'].value * 100
forward_Bolton_Index = ws2['P9'].value * 100

deviation_upper_canine_width = ws2['D23'].value
deviation_upper_premolars = ws2['D24'].value
deviation_upper_molars = ws2['D25'].value
length_upper_frontal_section = ws2['D26'].value

deviation_lower_canine_width = ws2['G23'].value
deviation_lower_premolars = ws2['G24'].value
deviation_lower_molars = ws2['G25'].value
length_lower_frontal_section = ws2['G26'].value

width_upper_dentition = ws2['C6'].value
width_lower_dentition = ws2['B7'].value
required_width_upper_dentition = ws2['E6'].value
required_width_lower_dentition = ws2['D7'].value
diff_width_upper_dentition = required_width_upper_dentition - width_upper_dentition

canine = 'клыками'
premolars = 'премолярами'
molars = 'молярами'

increased = 'увеличено'
decreased = 'уменьшено'

length_upper_frontal_status = ''
length_lower_frontal_status = ''

if length_upper_frontal_section > 1:
    length_upper_frontal_status = f"""Удлинение фронтального участка верхней челюсти на {format_with_comma(length_upper_frontal_section)} мм."""
elif length_upper_frontal_section < -1:
    length_upper_frontal_status = f"""Укорочение фронтального участка верхней челюсти на {process_string(length_upper_frontal_section, 1)} мм."""
else:
    length_upper_frontal_status = f"""Длина фронтального участка верхней челюсти в норме."""

if length_lower_frontal_section > 1:
    length_lower_frontal_status = f"""Удлинение фронтального участка нижней челюсти на {format_with_comma(length_lower_frontal_section)} мм."""
elif length_lower_frontal_section < -1:
    length_lower_frontal_status = f"""Укорочение фронтального участка нижней челюсти на {process_string(length_lower_frontal_section, 1)} мм."""
else:
    length_lower_frontal_status = f"""Длина фронтального участка нижней челюсти в норме."""

upper_frontal_without_number = ''
lower_frontal_without_number = ''

if length_upper_frontal_section > 1:
    upper_frontal_without_number = f"""Удлинение фронтального участка верхнего зубного ряда."""
elif length_upper_frontal_section < -1:
    upper_frontal_without_number = f"""Укорочение фронтального участка верхнего зубного ряда."""
else:
    upper_frontal_without_number = f"""Длина фронтального участка верхнего зубного ряда в норме."""

if length_lower_frontal_section > 1:
    lower_frontal_without_number = f"""Удлинение фронтального участка нижнего зубного ряда."""
elif length_lower_frontal_section < -1:
    lower_frontal_without_number = f"""Укорочение фронтального участка нижнего зубного ряда."""
else:
    lower_frontal_without_number = f"""Длина фронтального участка нижнего зубного ряда в норме."""


def analyze_tooth_type(deviation_upper_tooth_type_width, deviation_lower_tooth_type_width, tooth_type):
    tooth_type_result_str = ''
    if deviation_upper_tooth_type_width > 1 and deviation_lower_tooth_type_width > 1:
        tooth_type_result_str = f'Расстояние между {tooth_type} на верхней и нижней челюстях увеличено.'
    elif deviation_upper_tooth_type_width < -1 and deviation_lower_tooth_type_width < -1:
        tooth_type_result_str = f'Расстояние между {tooth_type} на верхней и нижней челюстях уменьшено.'
    elif deviation_upper_tooth_type_width > 1 and deviation_lower_tooth_type_width < -1:
        tooth_type_result_str = f'Расстояние между {tooth_type} на верхней челюсти увеличено, на нижней челюсти уменьшено.'
    elif deviation_upper_tooth_type_width < -1 and deviation_lower_tooth_type_width > 1:
        tooth_type_result_str = f'Расстояние между {tooth_type} на верхней челюсти уменьшено, на нижней челюсти увеличено.'
    elif 1 >= deviation_upper_tooth_type_width >= -1 > deviation_lower_tooth_type_width:
        tooth_type_result_str = f'Расстояние между {tooth_type} на верхней челюсти в норме, на нижней челюсти уменьшено.'
    elif -1 <= deviation_upper_tooth_type_width <= 1 < deviation_lower_tooth_type_width:
        tooth_type_result_str = f'Расстояние между {tooth_type} на верхней челюсти в норме, на нижней челюсти увеличено.'
    elif deviation_upper_tooth_type_width > 1 >= deviation_lower_tooth_type_width >= -1:
        tooth_type_result_str = f'Расстояние между {tooth_type} на верхней челюсти увеличено, на нижней челюсти в норме.'
    elif deviation_upper_tooth_type_width < -1 <= deviation_lower_tooth_type_width <= 1:
        tooth_type_result_str = f'Расстояние между {tooth_type} на верхней челюсти уменьшено, на нижней челюсти в норме.'
    elif -1 <= deviation_upper_tooth_type_width <= 1 and -1 <= deviation_lower_tooth_type_width <= 1:
        tooth_type_result_str = f'Расстояние между {tooth_type} на верхней и нижней челюстях в норме.'

    return tooth_type_result_str


mesiodystal_size = ""
if ton_index == 1.33:
    mesiodystal_size = 'пропорциональных'
else:
    mesiodystal_size = 'непропорциональных'

# Определяем класс в зависимости от значения ANB
anb_trend_class = ""
if anb_value > 4:
    anb_skeletal_class = "II"
elif anb_value < 0:
    anb_skeletal_class = "III"
else:
    anb_skeletal_class = "I"
    # Если класс "I", проверяем дополнительные условия
    if 0 <= anb_value <= 0.4:
        anb_trend_class = "с тенденцией к III классу"
    elif 3.6 <= anb_value <= 4:
        anb_trend_class = "с тенденцией к II классу"

# Определяем класс в зависимости от значения BETA ANGLE
beta_trend_class = ""
if beta_angle > 35:
    beta_skeletal_class = "III"
elif beta_angle < 27:
    beta_skeletal_class = "II"
else:
    beta_skeletal_class = "I"
    # Если класс "I", проверяем дополнительные условия
    if 34.6 <= beta_angle <= 35:
        beta_trend_class = "с тенденцией к III классу"
    elif 27 <= beta_angle <= 27.4:
        beta_trend_class = "с тенденцией к II классу"

# Определяем класс в зависимости от значения Wits Appraisal
has_value = ""
wits_trend_class = ""
if wits_appraisal > 2.1:
    wits_skeletal_class = "II"
    has_value = "наличие"
elif wits_appraisal < -2.9:
    wits_skeletal_class = "III"
    has_value = "наличие"
else:
    wits_skeletal_class = "I"
    has_value = "отсутствие"
    # Если класс "I", проверяем дополнительные условия
    if -2.5 <= wits_appraisal <= -2.9:
        wits_trend_class = "с тенденцией к III классу"
    elif 1.7 <= wits_appraisal <= 2.1:
        wits_trend_class = "с тенденцией к II классу"

# Определяем класс в зависимости от значения Wits Appraisal
sassouni_text = ""
sassouni_trend_class = ""
if sassouni > 3:
    sassouni_skeletal_class = "III"
elif sassouni < -3:
    sassouni_skeletal_class = "II"
else:
    sassouni_skeletal_class = "I"
    if 2.6 <= sassouni <= 3:
        sassouni_trend_class = "с тенденцией к III классу"
    elif -3 <= sassouni <= -2.6:
        sassouni_trend_class = "с тенденцией к II классу"

if sassouni > 0:
    has_direction = "кзади"
elif sassouni < 0:
    has_direction = "кпереди"

sassouni_text_not_null = f"""Соотношение челюстей по методике Sassouni говорит за {sassouni_skeletal_class} скелетный класс {sassouni_trend_class} — базальная дуга проходит на {process_string(sassouni, 2)} мм {has_direction} от точки В (N = 0,0 мм ± 3,0 мм)."""
sassouni_text_null = f"""Соотношение челюстей по методике Sassouni говорит за I скелетный класс — базальная дуга проходит через точку B (N = 0,0 мм ± 3,0 мм)."""

if sassouni == 0:
    sassouni_text = sassouni_text_null
else:
    sassouni_text = sassouni_text_not_null

# Определяем класс в зависимости от значения APDI
apdi_trend_class = ""
if apdi_value > 86.4:
    apdi_skeletal_class = "III"
elif apdi_value < 76.4:
    apdi_skeletal_class = "II"
else:
    apdi_skeletal_class = "I"
    # Если класс "I", проверяем дополнительные условия
    if 86 <= apdi_value <= 86.4:
        apdi_trend_class = "с тенденцией к III классу"
    elif 76.4 <= apdi_value <= 76.8:
        apdi_trend_class = "с тенденцией к II классу"

# Определяем класс в зависимости от значения PNSA
pnsa_upper_limit = ws1['D9'].value + 3.5
pnsa_lower_limit = ws1['D9'].value - 3.5
pnsa_trend_class = ""
pnsa_status = ""
if pnsa_value > pnsa_upper_limit:
    pnsa_status = "увеличению"
elif pnsa_value < pnsa_lower_limit:
    pnsa_status = "уменьшению"
else:
    pnsa_status = "норме"
    # Если класс "I", проверяем дополнительные условия
    if 86 <= pnsa_value <= 86.4:
        pnsa_trend_class = "с тенденцией к III классу"
    elif 76.4 <= pnsa_value <= 76.8:
        pnsa_trend_class = "с тенденцией к II классу"

# Определяем класс в зависимости от значения Ширины скелетного базиса (J-J)
jj_upper_limit = ws1['D10'].value + 3
jj_lower_limit = ws1['D10'].value - 3
jj_status = ""
if jj_value > jj_upper_limit:
    jj_status = "расширению"
elif jj_value < jj_lower_limit:
    jj_status = "сужению"
else:
    jj_status = "норме"

# Определяем класс в зависимости от значения SNA
if sna_value > 85:
    sna_status = "прогнати"
elif sna_value < 79:
    sna_status = "ретрогнати"
else:
    sna_status = "нормогнати"

sna_status_uppercase = sna_status.capitalize()

# Определяем класс в зависимости от значения PP/SN
if ppsn_value > 12:
    ppsn_status = "ретроинклинаци"
elif ppsn_value < 5:
    ppsn_status = "антеинклинаци"
else:
    ppsn_status = "нормоинклинаци"

ppsn_status_uppercase = ppsn_status.capitalize()


def check_disadvantage_width_upper_dentition(dentition_value):
    if dentition_value <= 0:
        return "Недостаток ширины зубного ряда на верхней челюсти отсутствует."
    else:
        return f"Недостаток ширины зубного ряда на верхней челюсти составляет {process_string(dentition_value, 1)} мм."


# условие для определения типа лицевой структуры на основе лицевого индекса
type_facial_structure = ''

if facial_index < 97:
    type_facial_structure = 'Брахиофациальный'
elif 97 <= facial_index <= 103:
    type_facial_structure = 'Мезофациальный'
elif facial_index > 103:
    type_facial_structure = 'Долихофациальный'

# Тут мы формируем логику для увеличения, уменьшения углов и генерации вывода
increased_angles = []
decreased_angles = []
normal_angles = []
middle_third_face = ws1['L69'].value

# Определение состояния для носолицевого угла
if nosocomial_angle < 36:
    decreased_angles.append('носолицевого угла')
elif 36 <= nosocomial_angle <= 38:
    normal_angles.append('носолицевого угла')
elif nosocomial_angle > 38:
    increased_angles.append('носолицевого угла')

# Определение состояния для носоподбородочного угла
if nasal_angle < 122:
    decreased_angles.append('носоподбородочного угла')
elif 122 <= nasal_angle <= 132:
    normal_angles.append('носоподбородочного угла')
elif nasal_angle > 132:
    increased_angles.append('носоподбородочного угла')

# Определение состояния для носогубного угла
if labial_angle < 100:
    decreased_angles.append('носогубного угла')
elif 100 <= labial_angle <= 110:
    normal_angles.append('носогубного угла')
elif labial_angle > 110:
    increased_angles.append('носогубного угла')

# Определение состояния для подбородочно-лицевого угла
if chin_facial_angle < 80:
    decreased_angles.append('подбородочно-лицевого угла')
elif 80 <= chin_facial_angle <= 90:
    normal_angles.append('подбородочно-лицевого угла')
elif chin_facial_angle > 90:
    increased_angles.append('подбородочно-лицевого угла')

# Словарь для нормальных состояний углов
angles_dict = {
    'носолицевого угла': 'носолицевой угол',
    'носоподбородочного угла': 'носоподбородочный угол',
    'носогубного угла': 'носогубный угол',
    'подбородочно-лицевого угла': 'подбородочно-лицевой угол'
}

# Формируем предложения
increased_sentence = ""
decreased_sentence = ""
normal_sentence = ""

# Формируем кадры для увеличенных углов
if increased_angles:
    increased_sentence = 'Увеличение ' + ', '.join(increased_angles) + '.'

# Формируем кадры для уменьшенных углов
if decreased_angles:
    decreased_sentence = 'Уменьшение ' + ', '.join(decreased_angles) + '.'

# Формируем кадры для нормальных углов
if normal_angles:
    normal_sentence = ' ' + ', '.join([angles_dict[angle] for angle in normal_angles]) + ' в норме.'

trimmed_message = normal_sentence.strip()  # Убираем пробелы для нормы

# Делаем первую букву заглавной
if trimmed_message:
    result_normal_message = trimmed_message[0].upper() + trimmed_message[1:]
else:
    result_normal_message = trimmed_message  # Если сообщение пустое, оставляем как есть

middle_third_face_status = ''

if middle_third_face < 47.5:
    middle_third_face_status = 'уменьшена'
elif 47.5 <= middle_third_face >= 52.5:
    middle_third_face_status = 'в норме'
elif middle_third_face > 52.5:
    middle_third_face_status = 'увеличена'

fotometrics_text = f"""
Угол выпуклости мягких тканей лица (gl-sn-pog) – {format_with_comma(soft_tissues_angle)}˚ (N = 163,0°-175,0°).
{type_facial_structure} тип строения лица.
Профиль лица {profile_type}.
{increased_sentence}
{decreased_sentence}
{result_normal_message}
Надподбородочная борозда не выражена. Средняя треть лица {middle_third_face_status}.
Положение губ относительно эстетической плоскости Ricketts: верхняя губа: {format_with_comma(upper_lip_position)} мм – {upper_lip_position_status},
нижняя губа:  {format_with_comma(lower_lip_position)} мм – {lower_lip_position_status} (N = -4,0 мм ± 2,0 мм.;  -2,0 мм ± 2,0 мм).
"""

# Формируем текст, вставляя значения переменных
biometrics_text = f"""
Окклюзия моляров по Энглю: справа III класс, слева III класс.
Окклюзия клыков по Энглю: справа III класс, слева III класс.
Индекс Тона = {process_string(ton_index, 2)}, что говорит о {mesiodystal_size} мезиодистальных размерах резцов на верхней и нижней челюсти (N = 1,33).
Общий Индекс Болтона = {process_string(general_Bolton_Index, 1)}% (N = 91,3%). Передний Индекс Болтона = {process_string(forward_Bolton_Index, 1)}% (N = 77,2%).
{analyze_tooth_type(deviation_upper_canine_width, deviation_lower_canine_width, canine)}
{analyze_tooth_type(deviation_upper_premolars, deviation_lower_premolars, premolars)}
{analyze_tooth_type(deviation_upper_molars, deviation_lower_molars, molars)}
{length_upper_frontal_status} {length_lower_frontal_status}
Глубина кривой Шпее справа – {process_string(depth_shpee_right_curve, 1)} мм, слева – {process_string(depth_shpee_left_curve, 1)} мм (N = 1,5 мм). Глубокая кривая Шпее справа \ и слева.
WALA Ridge анализ.
Ширина верхнего зубного ряда – {process_string(width_upper_dentition, 1)} мм, ширина нижнего зубного ряда – {process_string(width_lower_dentition, 1)} мм.
Требуемая ширина верхнего зубного ряда – {process_string(required_width_upper_dentition, 1)} мм. Требуемая ширина нижнего зубного ряда – {process_string(required_width_lower_dentition, 1)} мм.
{check_disadvantage_width_upper_dentition(diff_width_upper_dentition)}
"""

cephalometry_text = f"""
Межапикальный угол (<ANB) – {format_with_comma(anb_value)}˚, что соответствует соотношению челюстей по {anb_skeletal_class} скелетному классу {anb_trend_class} (N = 2,0˚ ± 2,0˚).
Угол Бета (< Beta Angle) – {format_with_comma(beta_angle)}˚, что cоответствует соотношению челюстей по {beta_skeletal_class} скелетному классу {beta_trend_class} (N = 31,0˚ ± 4,0˚).
Параметр Wits (Wits Appraisal.) –  {format_with_comma(wits_appraisal)} мм что указывает на {has_value} диспропорции в расположении апикальных базисов верхней и нижней челюстей в сагиттальной плоскости и говорит за {wits_skeletal_class} скелетный класс {wits_trend_class} (N = -1,1 мм ± 2,0 мм).
{sassouni_text}
Параметр APDI, указывающий на дисплазию развития челюстей в сагиттальной плоскости, равен {format_with_comma(apdi_value)}˚ и говорит за {apdi_skeletal_class} скелетный класс {apdi_trend_class} (N = 81,4˚ ± 5,0˚).
"""
resume_upper_jaw_text = f"""
Длина основания верхней челюсти (PNS-A) – {format_with_comma(pnsa_value)} мм, что соответствует {pnsa_status} (N = {format_with_comma(round(ws1['D9'].value, 1))} мм ± 3,5 мм).
Ширина основания верхней (J-J) челюсти –  {format_with_comma(jj_value)} мм, что соответствует {jj_status} (N = {format_with_comma(round(ws1['D10'].value, 1))} мм ± 3,0 мм):  справа – {format_with_comma(ws1['C11'].value)} мм, слева – {format_with_comma(ws1['C12'].value)} мм (N = {format_with_comma(ws1['D10'].value / 2)} мм ± 1,5 мм).
Положение верхней челюсти по сагиттали  (<SNA) – {format_with_comma(sna_value)}˚, что соответствует {sna_status}и (N = 82,0˚ ±  3,0˚).
Положение верхней челюсти по вертикали  (<SN-Palatal Plane) – {format_with_comma(ppsn_value)}˚, что соответствует {ppsn_status}и (N = 8,0˚ ± 3,0˚).
Roll ротация отсутствует\  вправо (по часовой стрелке) \ влево (против часовой стрелки).
Yaw ротация отсутствует \ вправо  (по часовой стрелке) \ влево (против часовой стрелки).
"""


def add_text_to_custom(prs_20, slide_index, left_20, top_20, width_20, height_20, text_20):
    slide = prs_20.slides[slide_index]
    textbox = slide.shapes.add_textbox(left_20, top_20, width_20, height_20).text_frame
    textbox.word_wrap = True
    paragraph_20 = textbox.add_paragraph()
    paragraph_20.font.size = Pt(10.5)
    paragraph_20.font.bold = False
    paragraph_20.font.name = "Montserrat"
    paragraph_20.text = text_20


# Использование функции с альтернативным названием
add_text_to_custom(prs, 20, Inches(0.4), Inches(0.8), Inches(7.21), Inches(2.5), fotometrics_text)
add_text_to_custom(prs, 20, Inches(0.4), Inches(3.15), Inches(7.21), Inches(2.5), biometrics_text)
add_text_to_custom(prs, 20, Inches(0.4), Inches(6.9), Inches(7.21), Inches(2.7), cephalometry_text)
add_text_to_custom(prs, 20, Inches(0.4), Inches(9.08), Inches(7.21), Inches(2.2), resume_upper_jaw_text)

print(f" Слайд 20 сформирован")

print("<-------------------------------------------------------------------------------------------------------->")
# Слайд 21
# Переменные с тенденциями к классу
go_me_r_value = ws1['C17'].value
go_me_l_value = ws1['C18'].value

go_go_r_value = ws1['C19'].value
go_go_l_value = ws1['C20'].value

ar_go_r_value = ws1['C21'].value
ar_go_l_value = ws1['C22'].value

md_md_value = ws1['c23'].value

snb_value = ws1['c24'].value

mp_sn_value = ws1['C25'].value

chin_displacement = ws1['C26'].value

ans_quotient = ws1['N9'].value

assessment_growth_type = ws1['L98'].value

ans_xi_pm = ws1['L113'].value

odi_value = ws1['L120'].value

u1_l1_r = ws1['C31'].value
u1_l1_l = ws1['C32'].value


def compare_value(value1, value2, name):
    if value1 is None or value2 is None:
        print(f"Ошибка: Недостаточно данных для сравнения {name}")
        return "None"
    return 'меньше' if value1 < value2 else 'больше'


# Определяем класс в зависимости от значения Go-Me
go_me_status = compare_value(go_me_r_value, go_me_l_value, "Go-Me")

# Определяем класс в зависимости от значения Go-Go
go_go_status = compare_value(go_go_r_value, go_go_l_value, "Go-Go")

# Определяем класс в зависимости от значения Ar-Go
ar_go_status = compare_value(ar_go_r_value, ar_go_l_value, "Ar-Go")

# Определяем класс в зависимости от значения Md-Md
md_upper_limit = ws1['D23'].value + 3
md_lower_limit = ws1['D23'].value - 3
md_status = ""
if md_md_value > md_upper_limit:
    md_status = "расширению"
elif md_md_value < md_lower_limit:
    md_status = "сужению"
else:
    md_status = "норме"

# Определяем класс в зависимости от значения <SNB
if snb_value > 83:
    snb_status = "прогнатии"
elif snb_value < 77:
    snb_status = "ретрогнати"
else:
    snb_status = "нормогнати"

# Определяем класс в зависимости от значения <MP\SN
if mp_sn_value > 36:
    mp_sn_status = "ретроинклинаци"
elif mp_sn_value < 28:
    mp_sn_status = "антеинклинаци"
else:
    mp_sn_status = "нормоинклинаци"

mp_sn_status_uppercase = mp_sn_status.capitalize()

# Определение смещения подбородка
if chin_displacement > 0:
    chin_displacement_status = f"влево на {process_string(chin_displacement, 1)} мм"
elif chin_displacement < 0:
    chin_displacement_status = f"вправо на {process_string(chin_displacement, 1)} мм"
else:
    chin_displacement_status = "не выявлено"

# Определяем класс в зависимости от значения (N-ANS) / (ANS-Gn)
if ans_quotient > 0.89:
    ans_quotient_status = "негармоничное"
elif ans_quotient < 0.71:
    ans_quotient_status = "негармоничное"
else:
    ans_quotient_status = "гармоничное"

# Определяем класс в зависимости от значения ANS-Xi-Pm
ans_xi_pm_upper_limit = round(ws1['N8'].value, 1) + 5.5
ans_xi_pm_lower_limit = round(ws1['N8'].value, 1) - 5.5

ans_xi_pm_status = ""
if ans_xi_pm > ans_xi_pm_upper_limit:
    ans_xi_pm_status = "увеличению"
elif ans_xi_pm < ans_xi_pm_lower_limit:
    ans_xi_pm_status = "уменьшению"
else:
    ans_xi_pm_status = "норме"

# Определяем класс в зависимости от значения ODI
if odi_value > 79.5:
    odi_value_status = "к глубокой резцовой окклюзии"
elif odi_value < 69.5:
    odi_value_status = "вертикальной резцовой дизокклюзии"
else:
    odi_value_status = "норме"


def get_tooth_status(slant_value, difference, upper_threshold, lower_threshold, tooth_num):
    if slant_value is None:
        return f"Нормальное положение зуба {tooth_num}"
    elif slant_value > upper_threshold:
        return f"Протрузия зуба  {tooth_num} на {process_string(difference, 1)}˚"
    elif slant_value < lower_threshold:
        return f"Ретрузия зуба  {tooth_num} на {process_string(difference, 1)}˚"
    else:
        return f"Нормальное положение зуба {tooth_num}"


def get_tooth_empty_status(slant_value, difference, upper_threshold, lower_threshold, tooth_num):
    if slant_value is None:
        return ""
    elif slant_value > upper_threshold:
        return f"Протрузия зуба  {tooth_num} на {process_string(difference, 1)}˚"
    elif slant_value < lower_threshold:
        return f"Ретрузия зуба  {tooth_num} на {process_string(difference, 1)}˚"
    else:
        return ""


slant_r1_1 = ws1['C33'].value
slant_l2_1 = ws1['C34'].value
slant_l3_1 = ws1['C36'].value
slant_r4_1 = ws1['C35'].value

slant_r1_1_dif = format_with_comma(round(ws1['G33'].value, 1)) if ws1['G33'].value is not None else None
slant_l2_1_dif = format_with_comma(round(ws1['G34'].value, 1)) if ws1['G34'].value is not None else None
slant_l3_1_dif = format_with_comma(round(ws1['G36'].value, 1)) if ws1['G36'].value is not None else None
slant_r4_1_dif = format_with_comma(round(ws1['G35'].value, 1)) if ws1['G35'].value is not None else None

r1_1_value_status = get_tooth_status(slant_r1_1, slant_r1_1_dif, 115, 105, 1.1)
l2_1_value_status = get_tooth_status(slant_l2_1, slant_l2_1_dif, 115, 105, 2.1)
l3_1_value_status = get_tooth_status(slant_l3_1, slant_l3_1_dif, 100, 90, 3.1)
r4_1_value_status = get_tooth_status(slant_r4_1, slant_r4_1_dif, 100, 90, 4.1)

r1_1_value_empty_status = get_tooth_empty_status(slant_r1_1, slant_r1_1_dif, 115, 105, 1.1)
l2_1_value_empty_status = get_tooth_empty_status(slant_l2_1, slant_l2_1_dif, 115, 105, 2.1)
l3_1_value_empty_status = get_tooth_empty_status(slant_l3_1, slant_l3_1_dif, 100, 90, 3.1)
r4_1_value_empty_status = get_tooth_empty_status(slant_r4_1, slant_r4_1_dif, 100, 90, 4.1)

result_string = ""

# Список статусов зубов
empty_status_list = [
    r1_1_value_empty_status,
    l2_1_value_empty_status,
    l3_1_value_empty_status,
    r4_1_value_empty_status
]

# Формируем строку вывода
for status in empty_status_list:
    if status:  # Если статус не пустой
        result_string += f"{status}. "

# Проверяем, если все статусы пустые
if not any(empty_status_list):  # Если в списке нет ни одного статуса
    result_string = "Наклон резцов в пределах нормы"

# Формирование строки с динамическими данными
u1_pp = f"{r1_1_value_status}. {l2_1_value_status}"
l1_pp = f"{l3_1_value_status}. {r4_1_value_status}"

# Формируем текст, вставляя значения переменных
resume_text2 = f"""
Длина тела нижней челюсти (Go-Me): справа – {format_with_comma(go_me_r_value)} мм, слева – {format_with_comma(go_me_l_value)}  мм (N = {format_with_comma(round(ws1['M59'].value, 1))} мм ± 5,0 мм).
Длина тела нижней челюсти справа {go_me_status}, чем слева на {format_with_comma(round(abs(go_me_r_value - go_me_l_value), 2))} мм.
Длина ветви нижней челюсти (Co-Go) : справа – {format_with_comma(go_go_r_value)}  мм,  слева – {format_with_comma(go_go_l_value)}  мм (N = {format_with_comma(round(ws1['D19'].value, 1))} мм ± 4,0 мм).
Длина ветви нижней челюсти справа {go_go_status}, чем слева на {format_with_comma(round(abs(go_go_r_value - go_go_l_value), 2))} мм.
Гониальный угол (<Ar-Go-Me): справа –  {format_with_comma(ar_go_r_value)}˚,  слева – {format_with_comma(ar_go_l_value)}˚ (N = {format_with_comma(round(ws1['D21'].value, 1))}˚ ± 5,0˚).
Гониальный угол справа {ar_go_status}, чем слева на {format_with_comma(round(abs(ar_go_r_value - ar_go_l_value), 2))}˚.
Ширина базиса нижней челюсти (Md-Md) – {format_with_comma(md_md_value)} мм, что соответствует {md_status} (N = {format_with_comma(round(ws1['D23'].value, 1))} мм ± 3,0 мм).
Положение нижней челюсти по сагиттали  (<SNB) – {format_with_comma(snb_value)}˚, что соответствует {snb_status}и (N = 80,0˚ ± 3,0˚).
Положение нижней челюсти по вертикали (<MP-SN) – {format_with_comma(mp_sn_value)}˚, что соответствует {mp_sn_status}и (N = 32,0˚ ± 4,0˚).
Смещение подбородка {chin_displacement_status}.
Roll ротация отсутствует \ вправо (по часовой стрелке) \ влево (против часовой стрелки).
Yaw ротация отсутствует \ вправо  (по часовой стрелке) \ влево (против часовой стрелки).
"""
resume_text3 = f"""
Вертикальное лицевое соотношение (N-ANS/ANS-Gn) {ans_quotient_status} – {round(ans_quotient, 2)} (N = 0,8 ± 0,09).
Отношение задней высоты лица к передней (S-Go/N-Gn) – {format_with_comma(assessment_growth_type)}% (N = 63,0% ± 2,0%).
Высота нижней трети лица по Ricketts (<ANS-Xi-Pm) – {format_with_comma(ans_xi_pm)}˚, что соответствует {ans_xi_pm_status} (N = IVP {format_with_comma(round(ws1['N8'].value, 1))}˚ ± 5,5˚).
Параметр ODI – {format_with_comma(odi_value)}˚, что соответствует {odi_value_status} (N = 74,5˚ ±  5,0˚).
"""

resume_text4 = f"""
Межрезцовый угол: справа – {format_with_comma(u1_l1_r)}˚, слева – {format_with_comma(u1_l1_l)}˚ (N = 130,0˚ ± 6,0˚).
{u1_pp} (N = 110,0˚± 5,0˚).
{l1_pp} (N = 95,0˚ ± 5,0˚).
"""

width_lower_jaw = ws1['C23'].value
width_upper_jaw = ws1['C10'].value
width_dif_jaw = width_lower_jaw + 5
jaw_dif = (width_upper_jaw - width_dif_jaw)

if jaw_dif < 0:
    jaw_status = f"составляет {format_with_comma(round(abs(jaw_dif), 1))} мм."
else:
    jaw_status = "отсутствует."

resume_text5 = f"""
Ширина базиса нижней челюсти – {format_with_comma(width_lower_jaw)} мм. Фактическая ширина базиса верхней челюсти – {format_with_comma(width_upper_jaw)} мм. 
Требуемая ширина базиса верхней челюсти = {format_with_comma(width_dif_jaw)} мм. 
Дефицит ширины скелетного базиса верхней челюсти {jaw_status}
"""
# Добавляем текст на слайд
text_width_5 = Inches(7.2)
text_height_5 = Inches(3.5)
text_left_5 = Inches(0.5)
text_top_5 = Inches(0.9)
name_textbox_5 = prs.slides[21].shapes.add_textbox(text_left_5, text_top_5, text_width_5, text_height_5)
text_frame = name_textbox_5.text_frame
text_frame.word_wrap = True
paragraph = text_frame.add_paragraph()
paragraph.font.size = Pt(10.5)
paragraph.font.bold = False
paragraph.font.name = "Montserrat"
paragraph.text = resume_text2

# Добавляем текст на слайд
text_left_21_2 = Inches(0.5)
text_top_21_2 = Inches(4.4)
name_textbox_21_2 = prs.slides[21].shapes.add_textbox(text_left_21_2, text_top_21_2, text_width_5, Inches(1))
text_frame = name_textbox_21_2.text_frame
text_frame.word_wrap = True
paragraph = text_frame.add_paragraph()
paragraph.font.size = Pt(10.5)
paragraph.font.bold = False
paragraph.font.name = "Montserrat"
paragraph.text = resume_text3

# Добавляем текст на слайд
text_left_21_3 = Inches(0.5)
text_top_21_3 = Inches(5.73)
name_textbox_21_3 = prs.slides[21].shapes.add_textbox(text_left_21_3, text_top_21_3, text_width_5, Inches(1))
text_frame = name_textbox_21_3.text_frame
text_frame.word_wrap = True
paragraph = text_frame.add_paragraph()
paragraph.font.size = Pt(10.5)
paragraph.font.bold = False
paragraph.font.name = "Montserrat"
paragraph.text = resume_text4

# Добавляем текст на слайд
name_textbox_21_4 = prs.slides[21].shapes.add_textbox(Inches(0.5), Inches(6.9), text_width_5, Inches(1))
text_frame = name_textbox_21_4.text_frame
text_frame.word_wrap = True
paragraph = text_frame.add_paragraph()
paragraph.font.size = Pt(10.5)
paragraph.font.bold = False
paragraph.font.name = "Montserrat"
paragraph.text = resume_text5

print(f" Слайд 21 сформирован")

print("<-------------------------------------------------------------------------------------------------------->")
# Слайд 22

pnsa_status_slide22 = ""
if pnsa_value > pnsa_upper_limit:
    pnsa_status_slide22 = "увеличен"
elif pnsa_value < pnsa_lower_limit:
    pnsa_status_slide22 = "уменьшен"
else:
    pnsa_status_slide22 = "в норме"

# Верхняя челюсть: пункт 1
if (width_dif_jaw - width_upper_jaw) <= 0:
    width_basis_lower_jaw = "Ширина базиса верхней челюсти в норме"
else:
    width_basis_lower_jaw = f"Сужение базиса верхней челюсти на {format_with_comma(round(abs(jaw_dif), 2))} мм"

# Нижняя челюсть: пункт 3
if go_me_r_value > ws1['D17'].value + ws1['E17'].value:
    go_me_r_status = "увеличена"
elif go_me_r_value < ws1['D17'].value - ws1['E17'].value:
    go_me_r_status = "уменьшена"
else:
    go_me_r_status = "в норме"

# Нижняя челюсть: пункт 3
if go_me_l_value > ws1['D18'].value + ws1['E18'].value:
    go_me_l_status = "увеличена"
elif go_me_l_value < ws1['D18'].value - ws1['E18'].value:
    go_me_l_status = "уменьшена"
else:
    go_me_l_status = "в норме"

# Нижняя челюсть: пункт 4
if go_go_r_value > ws1['D19'].value + ws1['E19'].value:
    go_go_r_status = "увеличена"
elif go_go_r_value < ws1['D19'].value - ws1['E19'].value:
    go_go_r_status = "уменьшена"
else:
    go_go_r_status = "в норме"

# Нижняя челюсть: пункт 4
# 59
if go_go_l_value > ws1['D20'].value + ws1['E20'].value:
    go_go_l_status = "увеличена"
elif go_go_l_value < ws1['D20'].value - ws1['E20'].value:
    go_go_l_status = "уменьшена"
else:
    go_go_l_status = "в норме"

snb_value_finish = ws1['c24'].value

# Определяем класс в зависимости от значения <SNB
if snb_value_finish > 83:
    snb_status_finish = "прогнатия"
elif snb_value_finish < 77:
    snb_status_finish = "ретрогнатия"
else:
    snb_status_finish = "нормогнатия"

snb_status_uppercase = snb_status.capitalize()

overbite_value = ws1['C29'].value
overjet_value = ws1['C30'].value

# Определяем класс в зависимости от значения Overbite
if overbite_value > 4.4:
    overbite_value_status = f"Глубокая резцовая окклюзия. Вертикальное резцовое перекрытие увеличено до {process_string(overbite_value, 1)} мм (N = 2,5 мм ± 2,0 мм)."
elif overbite_value < 0.5:
    overbite_value_status = f"Вертикальная резцовая дизокклюзия – {process_string(overbite_value, 1)} мм (N = 2,5 мм ± 2,0 мм)."
else:
    overbite_value_status = f"Вертикальное резцовое перекрытие в норме – {process_string(overbite_value, 1)} мм (N = 2,5 мм ± 2,0 мм)."

# Определяем класс в зависимости от значения Overjet
if overjet_value > 5:
    overjet_value_status = f"Сагиттальная щель – {process_string(overjet_value, 1)} мм (N = 2,5 мм ± 2,5 мм)."
elif overjet_value < 0:
    overjet_value_status = f"Обратная сагиттальная щель {process_string(overjet_value, 1)} мм (от -0,1 и выше) (N = 2,5 мм ± 2,5 мм)."
else:
    overjet_value_status = f"Сагиттальное резцовое перекрытие в норме – {process_string(overjet_value, 1)} мм (N = 2,5 мм ± 2,5 мм)."

# Определяем класс в зависимости от значения Md-Md
mdk_upper_limit = ws1['D23'].value + 3
mdk_lower_limit = ws1['D23'].value - 3
mdk_status = ""
if md_md_value > mdk_upper_limit:
    md_status = "Расширение базиса нижней челюсти относительно возрастной нормы."
elif md_md_value < mdk_lower_limit:
    md_status = "Сужение базиса нижней челюсти относительно возрастной нормы."
else:
    md_status = "Ширина базиса нижней челюсти в норме."


# Инициализация переменных для верхней и нижней челюстей
upper_jaw_displacement = ws1['M30'].value  # Смещение верхней челюсти
lower_jaw_displacement = ws1['N30'].value  # Смещение нижней челюсти

# Переменные для хранения описания смещений
upper_jaw_status = ''
lower_jaw_status = ''

# Обработка смещения верхней челюсти
if upper_jaw_displacement > 0:
    upper_jaw_status = f'Межрезцовая линия на верхней челюсти смещена относительно срединно \n    сагиттальной линии на {upper_jaw_displacement} мм вправо,'
elif upper_jaw_displacement < 0:
    upper_jaw_status = f'Межрезцовая линия на верхней челюсти смещена относительно срединно \n    сагиттальной линии на {abs(upper_jaw_displacement)} мм влево,'
else:
    upper_jaw_status = f'Межрезцовая линия на верхней челюсти не смещена относительно срединно \n    сагиттальной линии,'

# Обработка смещения нижней челюсти
if lower_jaw_displacement > 0:
    lower_jaw_status = f'на нижней челюсти смещена \n    относительно срединно сагиттальной линии на {lower_jaw_displacement} мм вправо.'
elif lower_jaw_displacement < 0:
    lower_jaw_status = f'на нижней челюсти смещена \n    относительно срединно сагиттальной линии на {abs(lower_jaw_displacement)} мм влево.'
else:
    lower_jaw_status = 'на нижней челюсти не смещена \n    относительно срединно сагиттальной линии.'

# Инициализация текстового фрейма на слайде
text_frame = prs.slides[22].shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(7), Inches(9)).text_frame
text_frame.word_wrap = True

# Общий текст
paragraph = text_frame.add_paragraph()
paragraph.font.size = Pt(11)
paragraph.font.bold = False
paragraph.font.name = "Montserrat"
paragraph.text = f"""
1. Скелетный III класс с тенденцией к III, обусловленный
   макро/микробнатией, про/ретрогнатией верхней/нижней челюсти,
   диспропорцией расположения апикальных базисов челюстей в сагиттальном
   направлении. Зубоальвеолярная форма дистальной/мезиальной окклюзии.
2. {type_facial_structure} тип строения лицевого отдела черепа.
3. Нейтральный тип роста с тенденцией к вертикальному/горизонтальному росту.
4. Высота нижней трети лица по Ricketts в {ans_xi_pm_status}.
5. Профиль лица {profile_type}.
6. Ретроположение верхней и нижней губы относительно эстетической
   плоскости Ricketts.
7. Сужение и уменьшение объема воздухоносных путей. Сужения и уменьшения
   объема воздухоносных путей не выявлено.
8. Нормальное/Переднее/Заднее положение правой/левой суставной головки
   височно-нижнечелюстного сустава.
9. Скелетный возраст соответствует IIIVI стадии созревания шейных позвонков.
"""

# Заголовок "Верхняя челюсть:"
upper_jaw_heading = text_frame.add_paragraph()
upper_jaw_heading.font.size = Pt(11)
upper_jaw_heading.font.bold = True
upper_jaw_heading.font.name = "Montserrat"
upper_jaw_heading.text = "Верхняя челюсть:"
upper_jaw_heading.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

# Основной текст для верхней челюсти
upper_jaw_content = text_frame.add_paragraph()
upper_jaw_content.font.size = Pt(11)
upper_jaw_content.font.bold = False
upper_jaw_content.font.name = "Montserrat"
upper_jaw_content.text = f"""
1. Размер основания верхней челюсти по сагиттали {pnsa_status_slide22}.
2. {width_basis_lower_jaw} (Penn анализ).
3. {sna_status_uppercase}я верхней челюсти. {ppsn_status_uppercase}я верхней челюсти.
4. Ротация верхней челюсти в Roll/Yaw плоскости вправо (по часовой стрелке) 
   \влево (против часовой стрелки).
"""

# Заголовок "Нижняя челюсть:"
lower_jaw_heading = text_frame.add_paragraph()
lower_jaw_heading.font.size = Pt(11)
lower_jaw_heading.font.bold = True
lower_jaw_heading.font.name = "Montserrat"
lower_jaw_heading.text = "Нижняя челюсть:"
lower_jaw_heading.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

# Основной текст для нижней челюсти
lower_jaw_content = text_frame.add_paragraph()
lower_jaw_content.font.size = Pt(11)
lower_jaw_content.font.bold = False
lower_jaw_content.font.name = "Montserrat"
lower_jaw_content.text = f"""
1. {md_status}
2. {snb_status_uppercase}я нижней челюсти. {mp_sn_status_uppercase}я нижней челюсти.
3. Длина тела нижней челюсти справа {go_me_r_status}. Длина тела нижней челюсти слева {go_me_l_status}.
4. Длина ветви нижней челюсти справа {go_go_r_status}. Длина ветви нижней челюсти слева {go_go_l_status}.
5. Смещение подбородка {chin_displacement_status} за счет скелетной асимметрии.
6. Ротация нижней челюсти в Roll/Yaw плоскости вправо (по часовой стрелке) 
   \влево (против часовой стрелки).
"""

# Заголовок "Параметры наклона и положения зубов:"
teeth_parameters_heading = text_frame.add_paragraph()
teeth_parameters_heading.font.size = Pt(11)
teeth_parameters_heading.font.bold = True
teeth_parameters_heading.font.name = "Montserrat"
teeth_parameters_heading.text = "Параметры наклона и положения зубов:"
teeth_parameters_heading.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

# Основной текст для параметров наклона и положения зубов
teeth_parameters_content = text_frame.add_paragraph()
teeth_parameters_content.font.size = Pt(11)
teeth_parameters_content.font.bold = False
teeth_parameters_content.font.name = "Montserrat"
teeth_parameters_content.text = f"""
1. {upper_jaw_status} {lower_jaw_status}
2. Сужение верхнего зубного ряда в области клыков, премоляров, моляров.
   Сужение нижнего зубного ряда в области клыков, моляров, премоляров.
3. {upper_frontal_without_number}
   {lower_frontal_without_number}
4. {result_string.strip()}.
5. {overbite_value_status}
6. {overjet_value_status}
7. Глубина кривой Шпее увеличена справа \ слева.
"""

# slide20_text1 = f"""
# 1. Скелетный III класс с тенденцией к III, обусловленный
#     макро \ микрогнатией, про \ ретрогнатией верхней \ нижней челюсти \
#     диспропорцией расположения апикальных базисов челюстей в сагиттальном
#     направлении.
#     Зубоальвеолярная форма дистальной \ мезиальной окклюзии.
# 2. {type_facial_structure} тип строения лицевого отдела черепа.
# 3. Нейтральный тип роста с тенденцией к вертикальному\ горизонтальному росту.
# 4. Высота нижней трети лица по Ricketts  в {ans_xi_pm_status}.
# 5. Профиль лица {profile_type}.
# 6. Ретроположение верхней и нижней губы относительно эстетической
#     плоскости Ricketts.
# 7. Сужение и уменьшение объема воздухоносных путей. Сужения и уменьшения
#     объема воздухоносных путей не выявлено.
# 8. Нормальное \ Переднее \ Заднее положение правой \ левой суставной головки
#     височно-нижнечелюстного сустава.
# 9. Скелетный возраст соответствует IIIVI стадии созревания шейных позвонков.
# \t\n
# Верхняя челюсть:
# 1. Размер основания верхней челюсти по сагиттали {pnsa_status_slide22}.
# 2. {width_basis_lower_jaw} (Penn анализ).
# 3. {sna_status_uppercase}я верхней челюсти. {ppsn_status_uppercase}я верхней челюсти.
# 4. Ротация верхней челюсти в Roll \Yaw плоскости вправо (по часовой стрелке)
#     \влево (против часовой стрелки).
# \t\n
# Нижняя челюсть:
# 1. {md_status}
# 2. {snb_status_uppercase}я нижней челюсти. {mp_sn_status_uppercase}я нижней челюсти.
# 3. Длина тела нижней челюсти справа {go_me_r_status}. Длина тела нижней челюсти слева {go_me_l_status}.
# 4. Длина ветви нижней челюсти справа {go_go_r_status}. Длина ветви нижней челюсти слева {go_go_l_status}.
# 5. Смещение подбородка {chin_displacement_status}, \ за счет скелетной асимметрии.
# 6. Ротация нижней челюсти в Roll \Yaw плоскости вправо (по часовой стрелке)
#     \влево (против часовой стрелки).
# \t\n
# Параметры наклона и положения зубов:
# 1. {upper_jaw_status} {lower_jaw_status}
# 2. Сужение верхнего зубного ряда в области клыков, премоляров, моляров.
#     Сужение нижнего зубного ряда в области клыков, моляров, премоляров.
# 3. {upper_frontal_without_number}
#     {lower_frontal_without_number}
# 4. {result_string.strip()}.
# 5. {overbite_value_status}
# 6. {overjet_value_status}
# 7. Глубина кривой Шпее в увеличена справа \ слева.
# """

# slide20_text4 = f"""
#
# """

#
# # Добавляем текст на слайд
# text_frame = prs.slides[22].shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(7), Inches(3.8)).text_frame
# text_frame.word_wrap = True
# paragraph = text_frame.add_paragraph()
# paragraph.font.size = Pt(11)
# paragraph.font.bold = False
# paragraph.font.name = "Montserrat"
# paragraph.text = slide20_text1
#
# # Добавляем текст на слайд
# text_frame = prs.slides[22].shapes.add_textbox(Inches(0.6), Inches(3.9), Inches(7), Inches(1.5)).text_frame
# text_frame.word_wrap = True
# paragraph = text_frame.add_paragraph()
# paragraph.font.size = Pt(11)
# paragraph.font.bold = False
# paragraph.font.name = "Montserrat"
# paragraph.text = slide20_text2
#
# # Добавляем текст на слайд
# text_frame = prs.slides[22].shapes.add_textbox(Inches(0.6), Inches(5.1), Inches(7), Inches(2.5)).text_frame
# text_frame.word_wrap = True
# paragraph = text_frame.add_paragraph()
# paragraph.font.size = Pt(11)
# paragraph.font.bold = False
# paragraph.font.name = "Montserrat"
# paragraph.text = slide20_text3
#
# # Добавляем текст на слайд
# text_frame = prs.slides[22].shapes.add_textbox(Inches(0.6), Inches(7.35), Inches(7), Inches(2.5)).text_frame
# text_frame.word_wrap = True
# paragraph = text_frame.add_paragraph()
# paragraph.font.size = Pt(11)
# paragraph.font.bold = False
# paragraph.font.name = "Montserrat"
# paragraph.text = slide20_text4

print(f" Слайд 22 сформирован")

print("<-------------------------------------------------------------------------------------------------------->")
if folder_name:
    save_folder = os.path.join(work_folder, folder_name)
    prs.save(os.path.join(save_folder, f"{folder_name}.pptx"))


def extract_text_from_slides(prs, slide_indices):
    """
    Функция для извлечения текста с указанных слайдов презентации.
    Args:
        prs (Presentation): Объект презентации PowerPoint.
        slide_indices (list): Список индексов слайдов для извлечения текста.
    Returns:
        str: Текст с указанных слайдов.
    """
    text = ""
    for idx in slide_indices:
        slide = prs.slides[idx]
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


# Индексы слайдов, с которых нужно извлечь текст
slide_indices_to_extract = [20, 21, 22]  # Пример: извлечение текста с первых трех слайдов

# Извлечение текста с указанных слайдов
extracted_text = extract_text_from_slides(prs, slide_indices_to_extract)

# Путь к файлу, в который будет сохранен текст
output_file_path = os.path.join(os.path.join(work_folder, folder_name), f"{folder_name}.txt")


def save_text_to_file(text, file_path):
    """
    Функция для сохранения текста в файл.
    Args:
        text (str): Текст для сохранения.
        file_path (str): Путь к файлу, в который будет сохранен текст.
    """
    with open(file_path, "w", encoding="utf-8") as file:
        file.write(text)


# Сохранение текста в файл
save_text_to_file(extracted_text, output_file_path)

print(f"😊 Текст успешно извлечен с выбранных слайдов и сохранен в файл {output_file_path}")
